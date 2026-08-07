#!/usr/bin/env python3
"""
Extract text from DOCX/PPTX into Markdown.
"""
import sys
from pathlib import Path
from docx import Document
from pptx import Presentation


def extract_docx(path: Path) -> str:
    doc = Document(str(path))
    lines = [f"# {path.name}\n"]
    for i, para in enumerate(doc.paragraphs, 1):
        text = para.text.strip()
        if text:
            lines.append(text)
    # tables
    for ti, table in enumerate(doc.tables, 1):
        lines.append(f"\n## Table {ti}")
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(lines)


def extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    lines = [f"# {path.name}\n"]
    for si, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Slide {si}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if shape.has_table:
                table = shape.table
                lines.append("\n[Table]")
                for row in table.rows:
                    cells = [cell.text.strip().replace("|", "\\|") for row_cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(lines)


def main():
    files = [
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/客货邮报告202602 v1.0.docx", "docx"),
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/关于平顶山市客货邮融合发展合作意向的汇报.docx", "docx"),
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/客货邮 位置码/开门红文件精神.docx", "docx"),
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/平顶山市客货邮融合发展合作意向汇报v1.1.pptx", "pptx"),
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/平顶山市客货邮融合（冷藏仓）汇报v1.2.pptx", "pptx"),
        ("/Users/yuelaiyuehao/Desktop/客货邮专题2026/平顶山市客货邮融合（冷藏仓）汇报v1.3.pptx", "pptx"),
        ("/Users/yuelaiyuehao/Desktop/汇报材料/平顶山市客货邮融合发展合作意向汇报（以平台技术展示为主）.pptx", "pptx"),
    ]
    out_dir = Path("/Users/yuelaiyuehao/Documents/workbuddy 设计创意空间/kehuoyou-knowledge/raw_extracts")
    out_dir.mkdir(parents=True, exist_ok=True)
    for fp, ftype in files:
        p = Path(fp)
        if not p.exists():
            print(f"SKIP (not found): {fp}")
            continue
        try:
            if ftype == "docx":
                md = extract_docx(p)
            else:
                md = extract_pptx(p)
            out_path = out_dir / f"{p.stem}.md"
            out_path.write_text(md, encoding="utf-8")
            print(f"EXTRACTED: {out_path}")
        except Exception as e:
            print(f"ERROR {fp}: {e}")


if __name__ == "__main__":
    main()
